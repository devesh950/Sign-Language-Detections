"""
Create beautiful PowerPoint presentation for Sign Language Detection project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BLUE = RGBColor(26, 35, 126)  # #1A237E
BLUE = RGBColor(33, 150, 243)      # #2196F3
LIGHT_BLUE = RGBColor(100, 181, 246)  # #64B5F6
GREEN = RGBColor(76, 175, 80)      # #4CAF50
ORANGE = RGBColor(255, 152, 0)     # #FF9800
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(66, 66, 66)

def add_title_slide(prs, title, subtitle):
    """Add title slide with gradient effect"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Emoji
    emoji_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(1), Inches(1), Inches(1)
    )
    emoji_frame = emoji_box.text_frame
    emoji_frame.text = "🤟"
    emoji_para = emoji_frame.paragraphs[0]
    emoji_para.font.size = Pt(72)
    emoji_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, points, bg_color=WHITE):
    """Add content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = bg_color
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(8.4), Inches(5)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.level = 0
        
    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points):
    """Add two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = WHITE
    background.line.fill.background()
    
    # Title bar
    title_bar = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE
    title_bar.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    
    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.8), Inches(4.2), Inches(5)
    )
    left_frame = left_box.text_frame
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.3), Inches(1.8), Inches(4.2), Inches(5)
    )
    right_frame = right_box.text_frame
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = GREEN
    
    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
    
    return slide

def add_stats_slide(prs):
    """Add statistics slide with numbers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(245, 245, 245)
    background.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "📊 Project Statistics"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Stats boxes
    stats = [
        ("37", "Gestures\nTrained", BLUE),
        ("3,700+", "Training\nImages", GREEN),
        ("21", "Hand\nLandmarks", ORANGE),
        ("63", "Input\nFeatures", RGBColor(156, 39, 176))
    ]
    
    x_start = 0.5
    width = 2.2
    
    for i, (number, label, color) in enumerate(stats):
        x = x_start + i * (width + 0.2)
        
        # Box
        box = slide.shapes.add_shape(
            1, Inches(x), Inches(2.2), Inches(width), Inches(3)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = color
        box.shadow.inherit = False
        
        # Number
        num_box = slide.shapes.add_textbox(
            Inches(x), Inches(2.5), Inches(width), Inches(1.2)
        )
        num_frame = num_box.text_frame
        num_frame.text = number
        num_para = num_frame.paragraphs[0]
        num_para.font.size = Pt(48)
        num_para.font.bold = True
        num_para.font.color.rgb = WHITE
        num_para.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(
            Inches(x), Inches(3.9), Inches(width), Inches(1)
        )
        label_frame = label_box.text_frame
        label_frame.text = label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(18)
        label_para.font.color.rgb = WHITE
        label_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_thank_you_slide(prs):
    """Add thank you slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background gradient
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()
    
    # Thank you
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Thank You! 🙏"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # GitHub
    github_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.5), Inches(8), Inches(0.8)
    )
    github_frame = github_box.text_frame
    github_frame.text = "GitHub: devesh950/Sign-Language-Detections"
    github_para = github_frame.paragraphs[0]
    github_para.font.size = Pt(24)
    github_para.font.color.rgb = LIGHT_BLUE
    github_para.alignment = PP_ALIGN.CENTER
    
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "Sign Language Detection", 
    "Real-time ASL Gesture Recognition using AI"
)

# Slide 2: Problem Statement
add_content_slide(prs, "❓ Problem Statement", [
    "• Over 70 million deaf people worldwide use sign language",
    "• Communication barrier between deaf and hearing communities",
    "• Limited availability of sign language interpreters",
    "• Need for accessible technology to bridge this gap",
    "• Real-time translation can improve daily interactions"
])

# Slide 3: Solution Overview
add_content_slide(prs, "💡 Our Solution", [
    "✅ Real-time ASL gesture recognition system",
    "✅ Uses computer vision and deep learning",
    "✅ Recognizes 37 gestures (A-Z, 0-9, Space)",
    "✅ Works with any webcam - no special hardware",
    "✅ Fast and accurate detection",
    "✅ Easy to deploy and use"
])

# Slide 4: Technology Stack
add_two_column_slide(prs, "🛠️ Technology Stack",
    "Computer Vision",
    [
        "OpenCV for video capture",
        "MediaPipe for hand detection",
        "21 hand landmark extraction",
        "Real-time processing"
    ],
    "Machine Learning",
    [
        "Keras/JAX neural network",
        "Dense layers architecture",
        "3,700+ training images",
        "High accuracy model"
    ]
)

# Slide 5: How It Works
add_content_slide(prs, "⚙️ How It Works", [
    "1️⃣  Capture video frame from webcam",
    "2️⃣  Detect hand using MediaPipe",
    "3️⃣  Extract 21 landmark points (63 coordinates)",
    "4️⃣  Feed landmarks to neural network",
    "5️⃣  Predict gesture with confidence score",
    "6️⃣  Display result in real-time"
])

# Slide 6: System Architecture
add_content_slide(prs, "🏗️ System Architecture", [
    "📹 Input Layer: Webcam video stream",
    "🔍 Processing: MediaPipe hand landmark detection",
    "🧠 Model: Dense Neural Network (63 → 128 → 64 → 37)",
    "📊 Output: Gesture prediction + confidence",
    "💻 Interface: OpenCV window / Streamlit web app"
])

# Slide 7: Features
add_content_slide(prs, "✨ Key Features", [
    "🎯 37 Gesture Recognition (A-Z, 0-9, Space)",
    "⚡ Real-time Processing (30+ FPS)",
    "🎤 Voice Output (Text-to-Speech)",
    "📺 Large On-Screen Subtitles",
    "🌐 Web Interface for Photo Upload",
    "📱 Works on Desktop and Laptop"
])

# Slide 8: Statistics
add_stats_slide(prs)

# Slide 9: Applications
add_two_column_slide(prs, "🎯 Applications",
    "Current Use Cases",
    [
        "Learning sign language",
        "Communication aid",
        "Accessibility tool",
        "Educational purposes"
    ],
    "Future Possibilities",
    [
        "Video call translation",
        "Public service kiosks",
        "Mobile applications",
        "Smart home control"
    ]
)

# Slide 10: Results & Performance
add_content_slide(prs, "📈 Results & Performance", [
    "✅ High Accuracy: 90%+ on test data",
    "✅ Fast Inference: <50ms per frame",
    "✅ Smooth Performance: 30+ FPS",
    "✅ Low Latency: Real-time detection",
    "✅ Robust: Works in various lighting conditions",
    "✅ Deployable: Runs on standard computers"
])

# Slide 11: Deployment Options
add_two_column_slide(prs, "🚀 Deployment",
    "Local Deployment ✅",
    [
        "Live camera detection",
        "Zero latency",
        "Python script",
        "Best performance"
    ],
    "Cloud Deployment ✅",
    [
        "Photo upload only",
        "Streamlit Cloud",
        "Accessible anywhere",
        "No installation"
    ]
)

# Slide 12: Thank You
add_thank_you_slide(prs)

# Save presentation
output_file = "Sign_Language_Detection_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created: {output_file}")
print(f"📊 Total slides: {len(prs.slides)}")
